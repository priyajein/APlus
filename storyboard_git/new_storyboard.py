import json
import asyncio
import websockets
from fastapi import FastAPI, WebSocket
from pptx import Presentation
from pptx.util import Inches
from crewai import Crew, Task, Agent, LLM
import groq


# Initialize FastAPI app
app = FastAPI()

# Groq AI Configuration (Replace with your actual API key)
GROQ_API_KEY = "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"
GROQ_MODEL = "llama-3.3-70b-versatile"

import litellm

# Example of initializing with Groq instead of OpenAI
llm = LLM(model="groq/llama-3.3-70b-versatile", api_key=GROQ_API_KEY)



def generate_slide_content(guidelines, slide_number):
    """Uses CrewAI with Groq Model to generate content for each slide."""
    
    # Get the slide title and keywords from the guidelines
    slide = guidelines['structure'][slide_number - 1]
    slide_title = slide['title']
    slide_keywords = ", ".join(slide['keywords'])
    
    # Get the topic from the guidelines
    topic = guidelines['topic']
    
    # Construct role, goal, and backstory based on the guidelines
    agent = Agent(
        name="Slide Content Generator",
        description=f"Generate content for Slide {slide_number}: {slide_title}",
        # model=GROQ_MODEL,
        # api_key=GROQ_API_KEY,
        role="content_generator",  # Role assigned to the agent
        goal=f"Generate a detailed, educational slide for '{slide_title}' under the topic '{topic}' based on the keywords: {slide_keywords}",
        backstory=f"Generate educational content about '{slide_title}' under the topic '{topic}', focusing on the keywords: {slide_keywords}. Use visuals like charts and diagrams if necessary.",
        llm = llm
    )
    
    # Construct the task using the guidelines
    task = Task(
        description=f"Create engaging and informative content for Slide {slide_number} titled '{slide_title}' under the topic '{topic}'",
        expected_output="Slide content in the form of bullet points with visuals.",  # Add expected_output field
        context=[  # Ensure context contains necessary fields
            {
                "description": f"Slide content for topic '{topic}' and slide '{slide_title}' with the keywords: {slide_keywords}.",
                "expected_output": "Bullet points with educational content and visuals."
            }
        ],
        agent=agent
        
    )
    

    # Initialize Crew and generate the slide content
    crew = Crew(
        agents = [agent],
        tasks=[task]
        )
    response = crew.kickoff()

    return response




def create_ppt(slides_content):
    """Creates a PowerPoint file from slide contents."""
    prs = Presentation()
    for idx, content in enumerate(slides_content, 1):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        title.text = f"Slide {idx}"
        body.text = content
    ppt_file = "storyboard.pptx"
    prs.save(ppt_file)
    return ppt_file

@app.websocket("/generate_ppt")
async def generate_ppt_endpoint(websocket: WebSocket):
    """WebSocket endpoint to stream slides as they are generated."""
    await websocket.accept()
    guidelines = await websocket.receive_text()
    with open("guidelines.json", "r") as file:
        guidelines_data = json.load(file)
    slides_content = []
    
    for slide_number in range(1, 21):
        slide_text = generate_slide_content(guidelines_data, slide_number)
        slides_content.append(slide_text)
        await websocket.send_json({"slide_number": slide_number, "content": slide_text})
    
    ppt_path = create_ppt(slides_content)
    await websocket.send_json({"status": "complete", "ppt_path": ppt_path})
    await websocket.close()

# Run the server using uvicorn
# Command: uvicorn filename:app --host 0.0.0.0 --port 8000